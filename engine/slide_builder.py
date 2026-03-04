"""
Slide builder — assembles .pptx files using python-pptx.
Builds branded Domo presentations from structured slide content.
"""

import copy
import io
from pathlib import Path
from typing import Dict, Any, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from config import BRAND, SLIDE_WIDTH_INCHES, SLIDE_HEIGHT_INCHES, ASSETS_DIR
from models.deck_config import LAYOUTS


# Brand colors as RGBColor objects
DOMO_BLUE = RGBColor(0x99, 0xCC, 0xEE)
DOMO_ORANGE = RGBColor(0xFF, 0x99, 0x22)
DARK_TEXT = RGBColor(0x3F, 0x45, 0x4D)
LIGHT_BG = RGBColor(0xF1, 0xF6, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x88, 0x88, 0x88)
BORDER = RGBColor(0xDC, 0xE4, 0xEA)
FONT_NAME = "Open Sans"

# Slide dimensions
SW = Inches(SLIDE_WIDTH_INCHES)
SH = Inches(SLIDE_HEIGHT_INCHES)


def build_presentation(
    slides: List[Dict[str, Any]],
    title: str = "Untitled",
    image_map: Optional[Dict[int, str]] = None,
    background_style: str = "white",
) -> bytes:
    """
    Build a complete .pptx from slide content.

    Args:
        slides: List of slide content dicts from content_writer
        title: Deck title (for metadata)
        image_map: Dict of {slide_position: image_path}
        background_style: "white" | "blue" | "gradient" | "dark"

    Returns:
        .pptx file as bytes
    """
    image_map = image_map or {}

    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    for slide_data in slides:
        layout_id = slide_data.get("layout_id", 12)
        layout_info = LAYOUTS.get(layout_id, {"type": "bullets"})
        layout_type = layout_info["type"]
        position = slide_data.get("position", 0)
        image_path = image_map.get(position)

        # Determine background based on layout
        is_blue = layout_id in (0, 3, 10, 22, 23, 35)
        is_gradient = layout_id in (2, 5)

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Set background
        if is_blue:
            _set_solid_bg(slide, DOMO_BLUE)
        elif is_gradient:
            _set_solid_bg(slide, RGBColor(0x5B, 0x89, 0xB3))
        # White layouts keep default white

        # Text colors based on background
        title_color = WHITE if (is_blue or is_gradient) else DARK_TEXT
        body_color = WHITE if (is_blue or is_gradient) else DARK_TEXT
        accent_color = DOMO_ORANGE if (is_blue or is_gradient) else DOMO_BLUE

        # Build slide based on type
        if layout_type == "title":
            _build_title_slide(slide, slide_data, title_color, accent_color)
        elif layout_type == "section":
            _build_section_slide(slide, slide_data, title_color, accent_color)
        elif layout_type == "quote":
            _build_quote_slide(slide, slide_data, title_color, accent_color)
        elif layout_type == "bullets":
            _build_bullets_slide(slide, slide_data, title_color, body_color, accent_color, image_path)
        elif layout_type == "one_column":
            _build_one_column_slide(slide, slide_data, title_color, body_color, accent_color)
        elif layout_type == "two_column":
            _build_two_column_slide(slide, slide_data, title_color, body_color, accent_color)
        elif layout_type == "text_image":
            _build_text_image_slide(slide, slide_data, title_color, body_color, accent_color, image_path)
        elif layout_type == "image":
            _build_image_slide(slide, slide_data, title_color, image_path)
        elif layout_type == "icons":
            _build_icons_slide(slide, slide_data, title_color, body_color, accent_color)
        elif layout_type == "close":
            _build_close_slide(slide, slide_data, title_color, accent_color)
        else:
            _build_bullets_slide(slide, slide_data, title_color, body_color, accent_color, image_path)

        # Add footer to all slides
        _add_footer(slide, position + 1, len(slides), is_blue or is_gradient)

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


# ── Background helpers ──────────────────────────────────────────────

def _set_solid_bg(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ── Slide type builders ─────────────────────────────────────────────

def _build_title_slide(slide, data, title_color, accent_color):
    headline = data.get("headline", "")
    subheadline = data.get("subheadline", "")

    # Logo placeholder area (top left)
    logo_path = ASSETS_DIR / "domo_logo_white.png"
    if not logo_path.exists():
        logo_path = ASSETS_DIR / "domo_logo_blue.png"
    if logo_path.exists():
        slide.shapes.add_picture(str(logo_path), Inches(0.5), Inches(0.4), height=Inches(0.45))

    # Title — large centered
    _add_text_box(slide, headline, Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.5),
                  font_size=Pt(44), bold=True, color=title_color, alignment=PP_ALIGN.LEFT)

    # Subtitle
    if subheadline:
        _add_text_box(slide, subheadline, Inches(0.8), Inches(3.8), Inches(9), Inches(0.8),
                      font_size=Pt(20), color=accent_color, alignment=PP_ALIGN.LEFT)

    # Accent line
    _add_line(slide, Inches(0.8), Inches(3.5), Inches(3), Pt(3), accent_color)


def _build_section_slide(slide, data, title_color, accent_color):
    headline = data.get("headline", "")
    subheadline = data.get("subheadline", "")

    _add_text_box(slide, headline.upper(), Inches(1), Inches(2.5), Inches(11), Inches(1.5),
                  font_size=Pt(40), bold=True, color=title_color, alignment=PP_ALIGN.CENTER)

    if subheadline:
        _add_text_box(slide, subheadline, Inches(2), Inches(4.2), Inches(9), Inches(0.7),
                      font_size=Pt(18), color=accent_color, alignment=PP_ALIGN.CENTER)


def _build_quote_slide(slide, data, title_color, accent_color):
    quote = data.get("quote", data.get("body", ""))
    attribution = data.get("quote_attribution", "")

    # Large quote mark
    _add_text_box(slide, "\u201C", Inches(1), Inches(1.5), Inches(2), Inches(1.5),
                  font_size=Pt(96), color=accent_color, bold=True)

    # Quote text
    _add_text_box(slide, quote, Inches(1.5), Inches(2.5), Inches(10), Inches(2.5),
                  font_size=Pt(24), color=title_color, italic=True)

    # Attribution
    if attribution:
        _add_text_box(slide, attribution, Inches(1.5), Inches(5.2), Inches(10), Inches(0.5),
                      font_size=Pt(14), color=accent_color)


def _build_bullets_slide(slide, data, title_color, body_color, accent_color, image_path=None):
    headline = data.get("headline", "")
    subheadline = data.get("subheadline", "")
    bullets = data.get("bullets", [])

    # Accent bar on left
    _add_rect(slide, Inches(0.4), Inches(0.4), Pt(4), Inches(0.8), accent_color)

    # Headline
    _add_text_box(slide, headline, Inches(0.7), Inches(0.35), Inches(11), Inches(0.75),
                  font_size=Pt(28), bold=True, color=title_color)

    if subheadline:
        _add_text_box(slide, subheadline, Inches(0.7), Inches(1.15), Inches(10), Inches(0.5),
                      font_size=Pt(14), color=GRAY)

    # Bullets
    content_width = Inches(7) if image_path else Inches(11)
    bullet_top = Inches(1.9) if subheadline else Inches(1.5)

    if bullets:
        tf_shape = slide.shapes.add_textbox(Inches(0.7), bullet_top, content_width, Inches(4.5))
        tf = tf_shape.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(8)
            run = p.add_run()
            run.text = f"\u2022  {bullet}"
            run.font.size = Pt(16)
            run.font.color.rgb = body_color
            run.font.name = FONT_NAME

    # Optional image on right
    if image_path and Path(image_path).exists():
        slide.shapes.add_picture(image_path, Inches(8.2), Inches(1.5), width=Inches(4.5))


def _build_one_column_slide(slide, data, title_color, body_color, accent_color):
    headline = data.get("headline", "")
    body = data.get("body", "")

    _add_rect(slide, Inches(0.4), Inches(0.4), Pt(4), Inches(0.8), accent_color)
    _add_text_box(slide, headline, Inches(0.7), Inches(0.35), Inches(11), Inches(0.75),
                  font_size=Pt(28), bold=True, color=title_color)

    if body:
        _add_text_box(slide, body, Inches(0.7), Inches(1.6), Inches(11.5), Inches(4.5),
                      font_size=Pt(16), color=body_color, line_spacing=Pt(28))


def _build_two_column_slide(slide, data, title_color, body_color, accent_color):
    headline = data.get("headline", "")
    bullets = data.get("bullets", [])
    body = data.get("body", "")

    _add_rect(slide, Inches(0.4), Inches(0.4), Pt(4), Inches(0.8), accent_color)
    _add_text_box(slide, headline, Inches(0.7), Inches(0.35), Inches(11), Inches(0.75),
                  font_size=Pt(28), bold=True, color=title_color)

    # Split bullets into two columns
    if bullets:
        mid = len(bullets) // 2 or 1
        left_bullets = bullets[:mid]
        right_bullets = bullets[mid:]

        _add_bullet_list(slide, left_bullets, Inches(0.7), Inches(1.6), Inches(5.5), Inches(4.5), body_color)
        _add_bullet_list(slide, right_bullets, Inches(6.8), Inches(1.6), Inches(5.5), Inches(4.5), body_color)

        # Vertical divider
        _add_rect(slide, Inches(6.55), Inches(1.6), Pt(1), Inches(4.5), BORDER)
    elif body:
        # Split body text
        sentences = body.split(". ")
        mid = len(sentences) // 2
        left_text = ". ".join(sentences[:mid]) + "."
        right_text = ". ".join(sentences[mid:])
        _add_text_box(slide, left_text, Inches(0.7), Inches(1.6), Inches(5.5), Inches(4.5),
                      font_size=Pt(14), color=body_color)
        _add_text_box(slide, right_text, Inches(6.8), Inches(1.6), Inches(5.5), Inches(4.5),
                      font_size=Pt(14), color=body_color)


def _build_text_image_slide(slide, data, title_color, body_color, accent_color, image_path=None):
    headline = data.get("headline", "")
    body = data.get("body", "")
    bullets = data.get("bullets", [])

    _add_rect(slide, Inches(0.4), Inches(0.4), Pt(4), Inches(0.8), accent_color)
    _add_text_box(slide, headline, Inches(0.7), Inches(0.35), Inches(6), Inches(0.75),
                  font_size=Pt(28), bold=True, color=title_color)

    if bullets:
        _add_bullet_list(slide, bullets, Inches(0.7), Inches(1.6), Inches(5.5), Inches(4.5), body_color)
    elif body:
        _add_text_box(slide, body, Inches(0.7), Inches(1.6), Inches(5.5), Inches(4.5),
                      font_size=Pt(14), color=body_color)

    if image_path and Path(image_path).exists():
        slide.shapes.add_picture(image_path, Inches(7), Inches(0.8), width=Inches(5.8))


def _build_image_slide(slide, data, title_color, image_path=None):
    headline = data.get("headline", "")

    if image_path and Path(image_path).exists():
        slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=SW)

    # Caption overlay at bottom
    if headline:
        _add_text_box(slide, headline, Inches(0.5), Inches(6.2), Inches(12), Inches(0.7),
                      font_size=Pt(18), color=title_color, bold=True)


def _build_icons_slide(slide, data, title_color, body_color, accent_color):
    headline = data.get("headline", "")
    icon_points = data.get("icon_points", [])

    _add_rect(slide, Inches(0.4), Inches(0.4), Pt(4), Inches(0.8), accent_color)
    _add_text_box(slide, headline, Inches(0.7), Inches(0.35), Inches(11), Inches(0.75),
                  font_size=Pt(28), bold=True, color=title_color)

    if not icon_points:
        return

    # Layout icons evenly across the slide
    count = min(len(icon_points), 3)
    col_width = 11 / count
    icon_symbols = ["\u25CF", "\u25B2", "\u25A0"]  # circle, triangle, square

    for i, point in enumerate(icon_points[:3]):
        x = Inches(0.7 + i * col_width)
        y_icon = Inches(2.2)
        y_label = Inches(3.5)
        y_desc = Inches(4.1)

        # Icon placeholder (colored circle)
        icon_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(col_width/2 - 0.4), y_icon, Inches(0.8), Inches(0.8)
        )
        icon_shape.fill.solid()
        icon_shape.fill.fore_color.rgb = accent_color
        icon_shape.line.fill.background()

        # Label
        label = point.get("label", f"Point {i+1}")
        _add_text_box(slide, label, x, y_label, Inches(col_width - 0.2), Inches(0.5),
                      font_size=Pt(16), bold=True, color=title_color, alignment=PP_ALIGN.CENTER)

        # Description
        desc = point.get("description", "")
        if desc:
            _add_text_box(slide, desc, x, y_desc, Inches(col_width - 0.2), Inches(1.5),
                          font_size=Pt(12), color=body_color, alignment=PP_ALIGN.CENTER)


def _build_close_slide(slide, data, title_color, accent_color):
    headline = data.get("headline", "Thank You")
    subheadline = data.get("subheadline", "")

    _add_text_box(slide, headline, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
                  font_size=Pt(44), bold=True, color=title_color, alignment=PP_ALIGN.CENTER)

    if subheadline:
        _add_text_box(slide, subheadline, Inches(2), Inches(4.2), Inches(9), Inches(0.7),
                      font_size=Pt(18), color=accent_color, alignment=PP_ALIGN.CENTER)

    # Accent line
    _add_line(slide, Inches(5.5), Inches(4.0), Inches(2.3), Pt(3), accent_color)


# ── Shared helpers ──────────────────────────────────────────────────

def _add_text_box(slide, text, left, top, width, height,
                  font_size=Pt(14), bold=False, italic=False,
                  color=DARK_TEXT, alignment=PP_ALIGN.LEFT, line_spacing=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = FONT_NAME
    return box


def _add_bullet_list(slide, items, left, top, width, height, color):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = f"\u2022  {item}"
        run.font.size = Pt(15)
        run.font.color.rgb = color
        run.font.name = FONT_NAME


def _add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_line(slide, left, top, width, height, color):
    return _add_rect(slide, left, top, width, height, color)


def _add_footer(slide, page_num, total_pages, light_text=False):
    text_color = WHITE if light_text else GRAY
    line_color = RGBColor(0xFF, 0xFF, 0xFF) if light_text else BORDER

    # Bottom line
    _add_rect(slide, Inches(0.4), Inches(6.9), Inches(12.5), Pt(1), line_color)

    # Logo text
    _add_text_box(slide, "DOMO", Inches(0.4), Inches(7.0), Inches(1.5), Inches(0.3),
                  font_size=Pt(9), bold=True, color=text_color)

    # Tagline
    _add_text_box(slide, "The AI and Data Products Platform", Inches(1.5), Inches(7.0), Inches(4), Inches(0.3),
                  font_size=Pt(8), color=text_color)

    # Page number
    _add_text_box(slide, f"{page_num}", Inches(12), Inches(7.0), Inches(1), Inches(0.3),
                  font_size=Pt(9), color=text_color, alignment=PP_ALIGN.RIGHT)

    # Confidential
    _add_text_box(slide, "Confidential", Inches(10), Inches(7.0), Inches(2), Inches(0.3),
                  font_size=Pt(8), color=text_color, alignment=PP_ALIGN.RIGHT)
